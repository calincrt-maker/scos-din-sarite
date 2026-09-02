import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- COLOR PALETTE ---
BG_COLOR = RGBColor(12, 13, 16)       # #0C0D10 Deep Charcoal/Black
CARD_BG = RGBColor(22, 25, 32)        # #161920 Elevated Card BG
CARD_BORDER = RGBColor(40, 45, 58)    # #282D3A Muted Card Border
RED_ACCENT = RGBColor(230, 20, 50)    # #E61432 Signature Neon Red
RED_LIGHT = RGBColor(255, 60, 85)     # #FF3C55 Soft Red Accent
WHITE = RGBColor(255, 255, 255)       # #FFFFFF Pure White
GRAY_LIGHT = RGBColor(200, 205, 215)  # #C8CDD7 Secondary Text
GRAY_MUTED = RGBColor(120, 125, 140)  # #787D8C Muted Text/Labels

FONT_MAIN = "Arial"
FONT_TITLE = "Trebuchet MS"

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def set_background(slide, color=BG_COLOR):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, category, title):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(9.5), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.name = FONT_MAIN
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = RED_ACCENT

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(9.5), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.name = FONT_TITLE
        p_title.font.size = Pt(23)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE

        # Intermediate Header Co-branding (New Digi24 Logo + HLM Minimalist Round Logo)
        digi_logo_path = "assets/digi24-dark-background.png" if os.path.exists("assets/digi24-dark-background.png") else "assets/logo_digi24_hd.png"
        if os.path.exists(digi_logo_path):
            slide.shapes.add_picture(digi_logo_path, Inches(11.0), Inches(0.4), Inches(0.65), Inches(0.65))
        if os.path.exists("assets/logo_hlm_round_minimal.png"):
            slide.shapes.add_picture("assets/logo_hlm_round_minimal.png", Inches(11.85), Inches(0.4), Inches(0.65), Inches(0.65))

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    # ==========================================
    # SLIDE 1: COVER (WITH 3 LOGOS TOP ROW)
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    set_background(s1)

    cover_img_path = "assets/cover_hero_full.jpg"
    if not os.path.exists(cover_img_path):
        cover_img_path = "assets/serghei_hero_studio.png"

    if os.path.exists(cover_img_path):
        s1.shapes.add_picture(cover_img_path, Inches(6.5), Inches(0.8), Inches(6.2), Inches(5.9))

    # 3 Production Logos on Slide 1 Top: Headline Management -> Digi24 -> Headline Production Hub
    if os.path.exists("assets/logo_headline_management.png"):
        s1.shapes.add_picture("assets/logo_headline_management.png", Inches(0.8), Inches(0.7), Inches(2.2), Inches(0.42))

    digi_logo_path = "assets/digi24-dark-background.png" if os.path.exists("assets/digi24-dark-background.png") else "assets/logo_digi24_hd.png"
    if os.path.exists(digi_logo_path):
        s1.shapes.add_picture(digi_logo_path, Inches(3.2), Inches(0.62), Inches(0.55), Inches(0.52))

    if os.path.exists("assets/logo_headline_production_hub.png"):
        s1.shapes.add_picture("assets/logo_headline_production_hub.png", Inches(3.9), Inches(0.7), Inches(2.2), Inches(0.42))

    # Main Title Box
    tb = s1.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.4), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "DIGI24  •  SEASON 2 PARTNERSHIP DECK"
    p.font.name = FONT_MAIN
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RED_ACCENT
    p.space_after = Pt(10)

    p2 = tf.add_paragraph()
    p2.text = "SCOS DIN SĂRITE\nCU SERGHEI"
    p2.font.name = FONT_TITLE
    p2.font.size = Pt(38)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_after = Pt(16)

    p3 = tf.add_paragraph()
    p3.text = "O perspectivă inteligentă și lucidă asupra cotidianului. Parteneriat comercial de nivel premium pe Digi24."
    p3.font.name = FONT_MAIN
    p3.font.size = Pt(13)
    p3.font.color.rgb = GRAY_LIGHT

    # ==========================================
    # SLIDE 2: COMBINED CONTEXT & CONCEPT EDITORIAL
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    set_background(s2)
    add_header(s2, "Context & Concept Editorial", "De la saturația media la o lentilă nouă: Umorul inteligent ca liant civic.")

    add_card(s2, Inches(0.8), Inches(1.8), Inches(11.733), Inches(0.9))
    tb_q = s2.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(11.333), Inches(0.7))
    tf_q = tb_q.text_frame
    tf_q.word_wrap = True
    p = tf_q.paragraphs[0]
    p.text = "„ Hazul de necaz reprezintă o formă de rezistență civică. Umorul nu evită realitatea — o face mai ușor de înțeles. ”"
    p.font.name = FONT_TITLE
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RED_LIGHT

    add_card(s2, Inches(0.8), Inches(2.9), Inches(5.6), Inches(4.0))
    tb_c1 = s2.shapes.add_textbox(Inches(1.0), Inches(3.1), Inches(5.2), Inches(3.6))
    tf_c1 = tb_c1.text_frame
    tf_c1.word_wrap = True

    p = tf_c1.paragraphs[0]
    p.text = "CLIMATUL MEDIA & NEVOIA AUDIENȚEI"
    p.font.name = FONT_MAIN
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RED_ACCENT
    p.space_after = Pt(8)

    p = tf_c1.add_paragraph()
    p.text = "• Newsfeed-ul zilnic și dezbaterile TV sunt marcate de conflict, rigiditate și polarizare obositoare.\n\n• Telespectatorii își doresc să rămână conectați, dar caută un spațiu de respiro: o supapă unde comedia de calitate detensionează cotidianul."
    p.font.name = FONT_MAIN
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY_LIGHT

    pillars = [
        ("01. ANALIZĂ PRIN COMEDIE", "Evenimentele săptămânii privite prin filtrul inteligent al observației sociale fin acordate."),
        ("02. ECHILIBRU FĂRĂ AGRESIVITATE", "Ironic dar decent; critic dar constructiv. O voce lucidă lipsită de vulgaritate sau atacuri."),
        ("03. INTELIGENT ȘI ACCESIBIL", "Umor contemporan ce vorbește pe înțelesul publicului urban, educat și conectat.")
    ]
    for i, (p_title, p_desc) in enumerate(pillars):
        top_pos = Inches(2.9 + i * 1.35)
        add_card(s2, Inches(6.6), top_pos, Inches(5.933), Inches(1.2))
        tb_p = s2.shapes.add_textbox(Inches(6.8), top_pos + Inches(0.15), Inches(5.533), Inches(0.9))
        tf_p = tb_p.text_frame
        tf_p.word_wrap = True

        p = tf_p.paragraphs[0]
        p.text = p_title
        p.font.name = FONT_MAIN
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_after = Pt(4)

        p = tf_p.add_paragraph()
        p.text = p_desc
        p.font.name = FONT_MAIN
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY_LIGHT

    # ==========================================
    # SLIDE 3: THE FORMAT
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    set_background(s3)
    add_header(s3, "Anatomia Formatului", "50 de minute dinamice de radiografie socială și comedie.")

    add_card(s3, Inches(0.8), Inches(2.0), Inches(6.5), Inches(4.8))
    tb_fmt = s3.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(5.9), Inches(4.2))
    tf_fmt = tb_fmt.text_frame
    tf_fmt.word_wrap = True

    features = [
        ("Current Affairs + Stand-up", "Combină rigoarea jurnalistică cu spontaneitatea comediei de observație."),
        ("Radiografie Socială", "De la derapajele din spațiul public la situațiile absurde din viața de zi cu zi."),
        ("Fără Tabuuri sau Prejudecăți", "Abordează subiecte fierbinți într-un cadru relaxat, optimist și decent."),
        ("Producție TV Contemporană", "Ritm alert, grafică modernă, decor de studio dinamic și atmosferă premium.")
    ]
    for i, (title, desc) in enumerate(features):
        p = tf_fmt.paragraphs[0] if i == 0 else tf_fmt.add_paragraph()
        p.text = f"• {title}"
        p.font.name = FONT_MAIN
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_after = Pt(2)

        p_sub = tf_fmt.add_paragraph()
        p_sub.text = desc
        p_sub.font.name = FONT_MAIN
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = GRAY_MUTED
        p_sub.space_after = Pt(10)

    img_s3 = "assets/serghei_portrait_v2.png" if os.path.exists("assets/serghei_portrait_v2.png") else "assets/serghei_studio_wide.png"
    if os.path.exists(img_s3):
        s3.shapes.add_picture(img_s3, Inches(7.6), Inches(2.0), Inches(3.6), Inches(4.8))

    # ==========================================
    # SLIDE 4: WHY SERGHEI
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    set_background(s4)
    add_header(s4, "Personality Spotlight", "Serghei — Vocea distinctă a noului val de divertisment.")

    img_serg = "assets/serghei_portrait_v1.png" if os.path.exists("assets/serghei_portrait_v1.png") else "assets/serghei_expressive_laugh.png"
    if os.path.exists(img_serg):
        s4.shapes.add_picture(img_serg, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.8))

    add_card(s4, Inches(4.8), Inches(2.0), Inches(7.733), Inches(4.8))
    tb_serg = s4.shapes.add_textbox(Inches(5.1), Inches(2.3), Inches(7.133), Inches(4.2))
    tf_serg = tb_serg.text_frame
    tf_serg.word_wrap = True

    p = tf_serg.paragraphs[0]
    p.text = "EDITORIAL PERSONALITY & HOST"
    p.font.name = FONT_MAIN
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RED_ACCENT
    p.space_after = Pt(10)

    p = tf_serg.add_paragraph()
    p.text = "Serghei aduce pe ecran o prezență autentică, construită pe observație fină, autoironie și dialog de calitate."
    p.font.name = FONT_TITLE
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(14)

    serg_traits = [
        "Capacitate dovedită de a detensiona subiecte delicate fără a deveni agresiv.",
        "Sute de mii de vizualizări în mediul online și spectacole live de succes.",
        "Discurs asumat, relaxat, lipsit de clișeele televiziunii clasice.",
        "Conexiune excelentă cu publicul tânăr-adult, educat și urban."
    ]
    for trait in serg_traits:
        p = tf_serg.add_paragraph()
        p.text = f"✓  {trait}"
        p.font.name = FONT_MAIN
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY_LIGHT
        p.space_after = Pt(8)

    # ==========================================
    # SLIDE 5: THE TONE
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    set_background(s5)
    add_header(s5, "Identitatea Tonală", "Trilovia calității: Cum diferențiem show-ul în peisajul TV.")

    tone_cards = [
        ("IRONIC", "NU CINIC", "Ironia taxează derapajul cu umor; cinismul îndepărtează audiența. Emisiunea păstrează o notă de optimism lucid."),
        ("CRITIC", "NU AGRESIV", "Critica aduce claritate și reflecție; agresivitatea creează zgomot. Serghei construiește dialog, nu atacuri."),
        ("INTELIGENT", "NU INACCESIBIL", "Umorul inteligent unește oamenii. Subiectele complexe sunt traduse simplu, firesc și memorabil.")
    ]
    for i, (main_t, sub_t, desc) in enumerate(tone_cards):
        left = Inches(0.8 + i * 4.0)
        add_card(s5, left, Inches(2.1), Inches(3.733), Inches(4.7))
        tb_t = s5.shapes.add_textbox(left + Inches(0.3), Inches(2.4), Inches(3.133), Inches(4.1))
        tf_t = tb_t.text_frame
        tf_t.word_wrap = True

        p = tf_t.paragraphs[0]
        p.text = main_t
        p.font.name = FONT_TITLE
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RED_LIGHT
        p.space_after = Pt(2)

        p = tf_t.add_paragraph()
        p.text = sub_t
        p.font.name = FONT_MAIN
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_after = Pt(18)

        p = tf_t.add_paragraph()
        p.text = desc
        p.font.name = FONT_MAIN
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY_LIGHT

    # ==========================================
    # SLIDE 6: SEASON 2 OVERVIEW
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    set_background(s6)
    add_header(s6, "Sezonul 2 în Cifre", "Parametrii de producție și difuzare ai Sezonului 2.")

    stats = [
        ("12", "EDIȚII SEZON II", "Producții săptămânale"),
        ("50 min", "DURATĂ SHOW", "Format dinamic fără timpi morți"),
        ("SEPT 2026", "LANSARE SEZON", "Grilă de toamnă Digi24"),
        ("3x", "EXPUNERE / EDIȚIE", "Premieră + 2 Redifuzări")
    ]
    for i, (num, label, sub) in enumerate(stats):
        left = Inches(0.8 + i * 3.0)
        add_card(s6, left, Inches(2.0), Inches(2.733), Inches(2.2))
        tb_s = s6.shapes.add_textbox(left + Inches(0.15), Inches(2.1), Inches(2.433), Inches(2.0))
        tf_s = tb_s.text_frame
        tf_s.word_wrap = True

        p = tf_s.paragraphs[0]
        p.text = num
        p.font.name = FONT_TITLE
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RED_LIGHT

        p = tf_s.add_paragraph()
        p.text = label
        p.font.name = FONT_MAIN
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_after = Pt(2)

        p = tf_s.add_paragraph()
        p.text = sub
        p.font.name = FONT_MAIN
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY_MUTED

    add_card(s6, Inches(0.8), Inches(4.5), Inches(11.733), Inches(2.3))
    tb_sch = s6.shapes.add_textbox(Inches(1.1), Inches(4.7), Inches(11.133), Inches(1.9))
    tf_sch = tb_sch.text_frame
    tf_sch.word_wrap = True

    p = tf_sch.paragraphs[0]
    p.text = "GRILA DE DIFUZARE DIGI24"
    p.font.name = FONT_MAIN
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RED_ACCENT
    p.space_after = Pt(10)

    schedules = [
        ("PREMIERĂ (PRIME TIME)", "Vineri: 20:00 – 21:00", "Interval de maximă audiență la debutul de weekend"),
        ("FIRST RE-RUN (EARLY PEAK)", "Sâmbătă: 19:00 – 20:00", "Expunere suplimentară în prime-time-ul de sâmbătă"),
        ("SECOND RE-RUN (LATE NIGHT)", "Sâmbătă / Duminică: 02:00 – 03:00", "Bonus inventory pentru publicul de noapte / nocturn")
    ]
    for title, time_slot, desc in schedules:
        p = tf_sch.add_paragraph()
        p.text = f"•  {title}  —  {time_slot}  |  {desc}"
        p.font.name = FONT_MAIN
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.space_after = Pt(6)

    # ==========================================
    # SLIDE 7: THE EXPERIENCE / GALLERY (SYMMETRIC WHITE BORDERS)
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    set_background(s7)
    add_header(s7, "Atmosfera din Studio", "Experiență vizuală și energie de producție premium.")

    images_g = [
        ("assets/serghei_portrait_v1.png", "Decizia editorială pe set"),
        ("assets/serghei_portrait_v2.png", "Interacțiunea din studio"),
        ("assets/cover_hero_full.jpg", "Recepția din Sezonul 1")
    ]
    for i, (img_path, caption) in enumerate(images_g):
        left = Inches(0.8 + i * 4.0)
        card = add_card(s7, left, Inches(2.0), Inches(3.733), Inches(4.8), bg_color=CARD_BG, border_color=WHITE)
        card.line.width = Pt(2)

        if os.path.exists(img_path):
            pic = s7.shapes.add_picture(img_path, left + Inches(0.15), Inches(2.15), Inches(3.433), Inches(3.6))
            pic.line.color.rgb = WHITE
            pic.line.width = Pt(1.5)

        tb_cap = s7.shapes.add_textbox(left + Inches(0.2), Inches(5.9), Inches(3.333), Inches(0.7))
        tf_cap = tb_cap.text_frame
        p = tf_cap.paragraphs[0]
        p.text = caption
        p.font.name = FONT_MAIN
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE

    # ==========================================
    # SLIDE 8: WHY IT WORKS FOR BRANDS
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    set_background(s8)
    add_header(s8, "Beneficii de Parteneriat", "De ce brandurile beneficiază de asocierea cu formatul.")

    brand_pillars = [
        ("RELEVANȚĂ CULTURALĂ", "Asociere directă cu conversația momentului și subiectele care preocupă publicul urban."),
        ("BRAND SAFETY GARANTAT", "Mediu TV premium pe Digi24. Umor inteligent, fără vulgaritate, atacuri ieftine sau risc reputațional."),
        ("RECEPȚIE POZITIVĂ", "Brandul apare într-un moment de relaxare și deschidere emoțională a telespectatorului."),
        ("CONTINUITATE & FRECVENȚĂ", "12 săptămâni consecutive de vizibilitate constantă pe TV și promovare în rețea.")
    ]
    for i, (b_title, b_desc) in enumerate(brand_pillars):
        col = i % 2
        row = i // 2
        left = Inches(0.8 + col * 6.0)
        top = Inches(2.0 + row * 2.5)
        add_card(s8, left, top, Inches(5.733), Inches(2.2))
        
        tb_b = s8.shapes.add_textbox(left + Inches(0.3), top + Inches(0.3), Inches(5.133), Inches(1.6))
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True

        p = tf_b.paragraphs[0]
        p.text = b_title
        p.font.name = FONT_TITLE
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RED_LIGHT
        p.space_after = Pt(8)

        p = tf_b.add_paragraph()
        p.text = b_desc
        p.font.name = FONT_MAIN
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY_LIGHT

    # ==========================================
    # SLIDE 9: PARTNERSHIP ECOSYSTEM
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    set_background(s9)
    add_header(s9, "Ecosistemul de Parteneriat", "Touchpoint-urile prin care brandul trăiește în emisiune.")

    ecosystem = [
        ("BBI & BBO", "Billboard In / Out", "Copertă de identificare a partenerului (5 sec) la începutul și finalul fiecărui bloc emisiune."),
        ("BREAKBUMPER", "Pauze Publicitare", "Marcaj vizual rapid (5 sec) la trecerea către și dinspre pauza de publicitate."),
        ("SQUEEZEBACK", "Animare In-Show", "Grafică animată pe ecran (10 sec) în timpul desfășurării emisiunii pentru memorabilitate."),
        ("SHARED PROMO", "Promovare TV Rețea", "Promovare TV dedicată emisiunii (5 sec) difuzată pe Digi24 pe tot parcursul săptămânii.")
    ]
    for i, (short_name, full_name, desc) in enumerate(ecosystem):
        left = Inches(0.8 + i * 3.0)
        add_card(s9, left, Inches(2.1), Inches(2.733), Inches(4.6))
        tb_e = s9.shapes.add_textbox(left + Inches(0.2), Inches(2.4), Inches(2.333), Inches(4.0))
        tf_e = tb_e.text_frame
        tf_e.word_wrap = True

        p = tf_e.paragraphs[0]
        p.text = short_name
        p.font.name = FONT_TITLE
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RED_LIGHT

        p = tf_e.add_paragraph()
        p.text = full_name
        p.font.name = FONT_MAIN
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_after = Pt(14)

        p = tf_e.add_paragraph()
        p.text = desc
        p.font.name = FONT_MAIN
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY_LIGHT

    # ==========================================
    # SLIDE 10: MEDIA EXPOSURE
    # ==========================================
    s10 = prs.slides.add_slide(blank_layout)
    set_background(s10)
    add_header(s10, "Planul de Expunere Media", "Sumarul complet al inventarului comercial pe cele 12 ediții.")

    add_card(s10, Inches(0.8), Inches(1.9), Inches(11.733), Inches(4.9))

    rows_data = [
        ("BBI & BBO (5s)", "2 / ed.", "24", "24", "24 (BONUS)", "72 Total"),
        ("Breakbumper (5s)", "4 / ed.", "48", "48", "48 (BONUS)", "144 Total"),
        ("Squeezeback (10s)", "2 / ed.", "24", "24", "24 (BONUS)", "72 Total"),
        ("Shared Promo (5s)", "35 / sapt.", "—", "—", "—", "420 Inserții TV")
    ]

    tb_t = s10.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(11.133), Inches(4.4))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True

    p = tf_t.paragraphs[0]
    p.text = f"{'ELEMENT COMERCIAL':<22} {'FRECVENȚĂ':<12} {'PREMIERĂ (VIE)':<16} {'RELUARE I (SÂM)':<16} {'RELUARE II (DUM)':<16} {'TOTAL SEZON'}"
    p.font.name = FONT_MAIN
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RED_ACCENT
    p.space_after = Pt(14)

    for el, freq, p_cnt, r1_cnt, r2_cnt, tot in rows_data:
        p = tf_t.add_paragraph()
        p.text = f"{el:<22} {freq:<12} {p_cnt:<16} {r1_cnt:<16} {r2_cnt:<16} {tot}"
        p.font.name = FONT_MAIN
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)

    p_summary = tf_t.add_paragraph()
    p_summary.text = "\nTOTAL EXPUNERE TV:  708 INSERȚII PE PARCURSUL A 12 SĂPTĂMÂNI"
    p_summary.font.name = FONT_TITLE
    p_summary.font.size = Pt(15)
    p_summary.font.bold = True
    p_summary.font.color.rgb = RED_LIGHT

    # ==========================================
    # SLIDE 11: COMMERCIAL OFFER
    # ==========================================
    s11 = prs.slides.add_slide(blank_layout)
    set_background(s11)
    add_header(s11, "Oferta Comercială", "Pachetul complet de sponsorizare și expunere media Sezonul 2.")

    add_card(s11, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.8))
    tb_pr = s11.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(5.0), Inches(4.2))
    tf_pr = tb_pr.text_frame
    tf_pr.word_wrap = True

    p = tf_pr.paragraphs[0]
    p.text = "VALOARE TOTALĂ PACHET"
    p.font.name = FONT_MAIN
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RED_ACCENT

    p = tf_pr.add_paragraph()
    p.text = "€ 51.000"
    p.font.name = FONT_TITLE
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(2)

    p = tf_pr.add_paragraph()
    p.text = "NET / 12 EDIȚII SEZONUL II"
    p.font.name = FONT_MAIN
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = GRAY_MUTED
    p.space_after = Pt(20)

    p = tf_pr.add_paragraph()
    p.text = "ECHIVALENT SĂPTĂMÂNAL:"
    p.font.name = FONT_MAIN
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = GRAY_LIGHT

    p = tf_pr.add_paragraph()
    p.text = "€ 4.250 NET / SĂPTĂMÂNĂ"
    p.font.name = FONT_TITLE
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RED_LIGHT

    add_card(s11, Inches(6.8), Inches(2.0), Inches(5.733), Inches(4.8))
    tb_inc = s11.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.133), Inches(4.2))
    tf_inc = tb_inc.text_frame
    tf_inc.word_wrap = True

    p = tf_inc.paragraphs[0]
    p.text = "PACHETUL INCLUDE INTEGRAL:"
    p.font.name = FONT_MAIN
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RED_ACCENT
    p.space_after = Pt(14)

    inclusions = [
        "708 Inserții TV totale pe parcursul celor 12 ediții.",
        "Exclusivitate pe categorie comercială (la nivel de sponsor principal).",
        "Pachet de 96 inserții BONUS incluse pe slotul de duminică.",
        "420 promo-uri TV de rețea difuzate pe Digi24.",
        "Cost mediu ultra-eficient per inserție TV: ~ € 72 / inserție."
    ]
    for inc in inclusions:
        p = tf_inc.add_paragraph()
        p.text = f"✓  {inc}"
        p.font.name = FONT_MAIN
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)

    # ==========================================
    # SLIDE 12: WHAT A BRAND BECOMES PART OF
    # ==========================================
    s12 = prs.slides.add_slide(blank_layout)
    set_background(s12)
    add_header(s12, "Dimensiunea Parteneriatului", "Mai mult decât spațiu publicitar: Asociere de impact.")

    img_s12 = "assets/cover_hero_full.jpg" if os.path.exists("assets/cover_hero_full.jpg") else "assets/serghei_hero_studio.png"
    if os.path.exists(img_s12):
        s12.shapes.add_picture(img_s12, Inches(6.6), Inches(1.8), Inches(5.933), Inches(4.8))

    add_card(s12, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.8))
    tb_q = s12.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(4.9), Inches(4.0))
    tf_q = tb_q.text_frame
    tf_q.word_wrap = True

    p = tf_q.paragraphs[0]
    p.text = "„ Brandurile memorabile nu întrerup conversația culturală. Ele devin parte din ea. ”"
    p.font.name = FONT_TITLE
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(20)

    p = tf_q.add_paragraph()
    p.text = "Scos din sărite cu Serghei oferă un spațiu de comunicare matur, inteligent și decent, unde brandul dumneavoastră este perceput ca un partener al bunei dispoziții și al umorului de calitate."
    p.font.name = FONT_MAIN
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY_LIGHT

    # ==========================================
    # SLIDE 13: CLOSING & MANDATORY PRODUCTION SECTION
    # Order: Headline Management -> Digi24 -> Headline Production Hub
    # Contact: Călin Crețu / calin@headliners.ro / +40 741 143 314
    # ==========================================
    s13 = prs.slides.add_slide(blank_layout)
    set_background(s13)

    add_card(s13, Inches(0.8), Inches(0.6), Inches(11.733), Inches(6.3))
    tb_cl = s13.shapes.add_textbox(Inches(1.3), Inches(0.8), Inches(10.733), Inches(2.0))
    tf_cl = tb_cl.text_frame
    tf_cl.word_wrap = True

    p = tf_cl.paragraphs[0]
    p.text = "SCOS DIN SĂRITE CU SERGHEI — SEZONUL 2"
    p.font.name = FONT_MAIN
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RED_ACCENT
    p.space_after = Pt(4)

    p = tf_cl.add_paragraph()
    p.text = "Construim împreună un parteneriat memorabil."
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(4)

    p = tf_cl.add_paragraph()
    p.text = "Haideți să discutăm integrarea brandului dumneavoastră în grila de toamnă 2026."
    p.font.name = FONT_MAIN
    p.font.size = Pt(13)
    p.font.color.rgb = RED_LIGHT

    # Production Credit Header: O PRODUCȚIE
    tb_prod = s13.shapes.add_textbox(Inches(1.3), Inches(2.9), Inches(10.733), Inches(0.4))
    tf_prod = tb_prod.text_frame
    p_p = tf_prod.paragraphs[0]
    p_p.alignment = PP_ALIGN.CENTER
    p_p.text = "O PRODUCȚIE"
    p_p.font.name = FONT_MAIN
    p_p.font.size = Pt(11)
    p_p.font.bold = True
    p_p.font.color.rgb = GRAY_MUTED

    # Mandatory 3-Logo Row: 1. Headline Management | 2. Digi24 | 3. Headline Production Hub
    if os.path.exists("assets/logo_headline_management.png"):
        s13.shapes.add_picture("assets/logo_headline_management.png", Inches(1.8), Inches(3.4), Inches(3.4), Inches(0.6))
    
    digi_logo_path = "assets/digi24-dark-background.png" if os.path.exists("assets/digi24-dark-background.png") else "assets/logo_digi24_hd.png"
    if os.path.exists(digi_logo_path):
        s13.shapes.add_picture(digi_logo_path, Inches(6.2), Inches(3.3), Inches(0.9), Inches(0.8))
        
    if os.path.exists("assets/logo_headline_production_hub.png"):
        s13.shapes.add_picture("assets/logo_headline_production_hub.png", Inches(8.1), Inches(3.4), Inches(3.4), Inches(0.6))

    # Contact Block: Călin Crețu / calin@headliners.ro / +40 741 143 314
    tb_cnt = s13.shapes.add_textbox(Inches(1.3), Inches(4.5), Inches(10.733), Inches(1.8))
    tf_cnt = tb_cnt.text_frame
    tf_cnt.word_wrap = True

    p = tf_cnt.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Călin Crețu"
    p.font.name = FONT_TITLE
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(6)

    p = tf_cnt.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.text = "calin@headliners.ro   |   +40 741 143 314"
    p.font.name = FONT_MAIN
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RED_LIGHT

    output_path = "Scos_din_sarite_cu_Serghei_Sezonul_2_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_deck()
